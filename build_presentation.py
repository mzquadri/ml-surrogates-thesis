"""Build 15-slide PowerPoint deck for the thesis defence.
Light pastel TUM palette, embedded figures, speaker notes attached.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ─────────────────────────── Paths ───────────────────────────
BASE = r"C:\Users\zamin\Downloads\ml_surrogates_thesis_final\ml_surrogates_thesis_final\document\figures\new"
OUT_PATH = r"C:\Users\zamin\Downloads\ml_surrogates_thesis_final\Zamin_Thesis_Defence_15_Slides.pptx"

# ────────────────────── Light TUM palette ────────────────────
TUM_BLUE_DARK   = RGBColor(0x00, 0x52, 0x93)   # for titles
TUM_GREY_LIGHT  = RGBColor(0xDA, 0xD7, 0xCB)
TUM_BLUE_LIGHT  = RGBColor(0x98, 0xC6, 0xEA)
TUM_BLUE_PALE   = RGBColor(0xCD, 0xE5, 0xF2)
TUM_ORANGE      = RGBColor(0xF4, 0xB1, 0x83)
TUM_CORAL       = RGBColor(0xE8, 0xA5, 0xA5)
TEXT_DARK       = RGBColor(0x22, 0x22, 0x22)
TEXT_GREY       = RGBColor(0x55, 0x55, 0x55)
ACCENT_LINE     = RGBColor(0x88, 0x88, 0x88)

# ─────────────────────── Setup deck ──────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]

# ─────────────────────── Helpers ─────────────────────────────


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, italic=False,
             color=TEXT_DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top  = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return box


def add_bullets(slide, x, y, w, h, bullets, *, size=14, color=TEXT_DARK,
                line_spacing=1.18, font="Calibri"):
    """bullets: list of strings; lines starting with '  ' (2 spaces) are sub-bullets."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top  = Emu(0); tf.margin_bottom = Emu(0)
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        # Sub-bullet detection
        if line.startswith("  "):
            p.level = 1
            text_to_show = "• " + line.strip()
            sz = size - 2
        elif line.startswith("→") or line.startswith("✅") or line.startswith("⚠"):
            p.level = 0
            text_to_show = line
            sz = size
        elif line == "":
            p.level = 0
            text_to_show = " "
            sz = size
        else:
            p.level = 0
            text_to_show = line
            sz = size
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text_to_show
        run.font.size = Pt(sz)
        run.font.color.rgb = color
        run.font.name = font
    return box


def add_title_bar(slide, slide_num, title_text):
    """Top title bar with slide number badge and title."""
    # thin coloured bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = TUM_BLUE_DARK
    bar.line.fill.background()
    # slide number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.32),
                                    Inches(0.7), Inches(0.45))
    badge.fill.solid(); badge.fill.fore_color.rgb = TUM_BLUE_LIGHT
    badge.line.fill.background()
    badge_tf = badge.text_frame
    badge_tf.margin_left = Emu(0); badge_tf.margin_right = Emu(0)
    badge_p = badge_tf.paragraphs[0]
    badge_p.alignment = PP_ALIGN.CENTER
    badge_run = badge_p.add_run()
    badge_run.text = f"{slide_num}"
    badge_run.font.size = Pt(20)
    badge_run.font.bold = True
    badge_run.font.color.rgb = TUM_BLUE_DARK
    badge_run.font.name = "Calibri"
    # title text
    add_text(slide, Inches(1.25), Inches(0.32), Inches(11.5), Inches(0.55),
             title_text, size=24, bold=True, color=TUM_BLUE_DARK)


def add_footer(slide, slide_num):
    add_text(slide, Inches(0.4), Inches(7.15), Inches(12), Inches(0.3),
             f"M.Zamin Quadri  ·  Master's Thesis Defence  ·  TU Munich  ·  May 2026",
             size=9, color=TEXT_GREY, italic=True)
    add_text(slide, Inches(12.55), Inches(7.15), Inches(0.6), Inches(0.3),
             f"{slide_num} / 15", size=9, color=TEXT_GREY, italic=True)


def add_image_safe(slide, x, y, w, h, filename):
    """Insert image if file exists, else placeholder."""
    full = os.path.join(BASE, filename)
    if os.path.isfile(full):
        slide.shapes.add_picture(full, x, y, width=w, height=h)
    else:
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        ph.fill.solid(); ph.fill.fore_color.rgb = TUM_GREY_LIGHT
        ph.line.color.rgb = ACCENT_LINE
        tf = ph.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = f"[figure: {filename}]"
        run.font.size = Pt(11); run.font.italic = True
        run.font.color.rgb = TEXT_GREY


def set_notes(slide, text):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


# ──────────────────── SLIDE 1 — Title ────────────────────────
s = prs.slides.add_slide(BLANK)

# accent bar at top
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.25))
bar.fill.solid(); bar.fill.fore_color.rgb = TUM_BLUE_DARK
bar.line.fill.background()

add_text(s, Inches(0.8), Inches(1.0), Inches(11.7), Inches(1.4),
         "Uncertainty Quantification for Machine Learning Models",
         size=34, bold=True, color=TUM_BLUE_DARK, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(1.0),
         "in Transportation Policy Analysis",
         size=34, bold=True, color=TUM_BLUE_DARK, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.8), Inches(3.0), Inches(11.7), Inches(0.5),
         "Mohd Zamin Quadri",
         size=22, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(3.5), Inches(11.7), Inches(0.4),
         "Master's Thesis · Technical University of Munich",
         size=16, italic=True, color=TEXT_GREY, align=PP_ALIGN.CENTER)

add_text(s, Inches(3.5), Inches(4.4), Inches(6.5), Inches(0.4),
         "Examiner    :  Prof. Dr. Stephan Günnemann",
         size=14, color=TEXT_DARK, font="Consolas")
add_text(s, Inches(3.5), Inches(4.78), Inches(6.5), Inches(0.4),
         "Supervisors :  Dominik Fuchsgruber  ·  Elena Natterer",
         size=14, color=TEXT_DARK, font="Consolas")

# message strip
strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(5.6),
                           Inches(10.3), Inches(1.0))
strip.fill.solid(); strip.fill.fore_color.rgb = TUM_BLUE_PALE
strip.line.fill.background()
tf = strip.text_frame
tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
tf.margin_top = Inches(0.15); tf.margin_bottom = Inches(0.15)
tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = ("Adding per-prediction uncertainty to a PointNetTransfGAT traffic surrogate "
          "for capacity-reduction scenarios on the Paris road network.")
r.font.size = Pt(15); r.font.italic = True; r.font.color.rgb = TUM_BLUE_DARK
r.font.name = "Calibri"

add_text(s, Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4),
         "May 2026",
         size=14, italic=True, color=TEXT_GREY, align=PP_ALIGN.CENTER)

set_notes(s, """Good afternoon everyone, and thank you for being here. My thesis studies uncertainty quantification for machine-learning surrogates in transport policy analysis. The starting point is the PointNetTransfGAT surrogate that Elena developed for the Paris road network. That surrogate is useful because a single MATSim simulation takes hours, while the GNN can produce per-segment predictions in seconds. But for policy analysis, a fast prediction is not enough on its own. A planner also needs to know when the surrogate is reliable, and when a prediction should be treated with caution or rerun in MATSim. So the focus of my thesis is to add per-prediction uncertainty estimates around this surrogate, evaluate how trustworthy they are, and be clear about where their uncertainty estimates remain limited. Everything I show today is empirical, and scoped to one Paris network, one 1,000-scenario subset out of 10,000, one type of policy intervention which is capacity reduction, and mainly one architecture family, PointNetTransfGAT.""")


# ──────────────────── SLIDE 2 — Agenda ───────────────────────
s = prs.slides.add_slide(BLANK)
add_title_bar(s, 2, "Agenda")
add_bullets(s, Inches(1.0), Inches(1.5), Inches(11.3), Inches(5.5), [
    "1.  Why traffic surrogates need uncertainty",
    "",
    "2.  Dataset, prediction task, and the PointNetTransfGAT surrogate",
    "",
    "3.  Base trials T1–T8 and the T8 UQ checkpoint",
    "",
    "4.  UQ methods: MC Dropout, calibration, conformal, ensembles",
    "",
    "5.  Uncertainty-aware training: T9, T10, T11",
    "",
    "6.  Non-GNN baselines, limitations, and takeaways",
], size=20)
add_footer(s, 2)
set_notes(s, "Before I go into the details, let me quickly outline how the talk is structured. I begin with the motivation. Then I describe the prediction task and the PointNetTransfGAT surrogate that this thesis builds on, which is Elena's architecture. After that I walk through the base training trials and explain why Trial 8 became the main UQ checkpoint. The central part is the UQ analysis itself: MC Dropout, temperature scaling, conformal prediction, and ensembles. I cover the three uncertainty-aware training variants — T9, T10, T11 — and the non-GNN baselines that put the GNN's point accuracy in context. I finish with main findings, limitations, and takeaways.")


# ──────────────────── SLIDE 3 — Motivation ───────────────────
s = prs.slides.add_slide(BLANK)
add_title_bar(s, 3, "Motivation: MATSim is accurate but slow")
add_bullets(s, Inches(0.9), Inches(1.4), Inches(11.5), Inches(4), [
    "MATSim — the standard agent-based simulator for transport policy",
    "  Models every traveller's daily plan (activities, modes, routes)",
    "  Iterates to a user equilibrium → high fidelity",
    "  One Paris-scale scenario takes more than 8 hours",
    "",
    "A planner comparing 12 candidate policies → ≈ 100 simulation hours",
], size=18)
# highlight box
hl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(5.3),
                        Inches(11.5), Inches(1.3))
hl.fill.solid(); hl.fill.fore_color.rgb = TUM_BLUE_PALE
hl.line.color.rgb = ACCENT_LINE
tf = hl.text_frame; tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2)
tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Fast surrogates make this practical — only if their predictions can be trusted."
r.font.size = Pt(18); r.font.italic = True; r.font.color.rgb = TUM_BLUE_DARK
r.font.bold = True; r.font.name = "Calibri"
add_footer(s, 3)
set_notes(s, "MATSim is the standard simulator for transport policy work. It models every traveller individually and iterates to user equilibrium. That detail makes it trustworthy, but slow — a single Paris scenario takes more than eight hours. A planner comparing twelve candidate policies needs close to a hundred hours of compute. Elena's GNN surrogate cuts this from hours to seconds. But seconds without reliability is a step backwards — a fast wrong answer is worse than a slow correct one. That's where uncertainty quantification comes in.")


# ──────────────────── SLIDE 4 — Problem formulation ──────────
s = prs.slides.add_slide(BLANK)
add_title_bar(s, 4, "Problem formulation: predict per-segment Δv")

add_bullets(s, Inches(0.9), Inches(1.4), Inches(7.0), Inches(5), [
    "Each Paris scenario → directed line graph",
    "  31,635 road segments (= nodes)",
    "",
    "Input per segment (5 features):",
    "  VOL_BASE_CASE",
    "  CAPACITY_BASE_CASE",
    "  CAPACITY_REDUCTION",
    "  FREESPEED",
    "  LENGTH",
    "",
    "Target: Δv — change in vehicles/hour",
    "",
    "Regression at the node level",
    "1,000 scenarios → 80 / 10 / 10 split",
    "→ ≈ 3.16 M test predictions",
], size=15)

# numbers panel on the right
panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.4), Inches(1.5),
                           Inches(4.4), Inches(5.0))
panel.fill.solid(); panel.fill.fore_color.rgb = TUM_BLUE_PALE
panel.line.color.rgb = ACCENT_LINE
add_text(s, Inches(8.7), Inches(1.65), Inches(4.0), Inches(0.4),
         "Key numbers", size=16, bold=True, color=TUM_BLUE_DARK)
add_bullets(s, Inches(8.7), Inches(2.15), Inches(4.0), Inches(4.2), [
    "31,635 nodes / scenario",
    "",
    "1,000 scenarios used",
    "(of 10,000 available)",
    "",
    "Train: 800 graphs",
    "Val:   100 graphs",
    "Test:  100 graphs",
    "",
    "= 3,163,500 test predictions",
], size=14, color=TEXT_DARK)
add_footer(s, 4)
set_notes(s, "The prediction task is at the level of individual road segments. Each Paris scenario is a directed line graph — every road segment becomes a node. The full Paris network has just over 31,000 segments. For every node we have five input features describing the segment and the policy. The target is delta v — the change in traffic volume induced by the policy. From the 10,000 scenarios that Elena's pipeline can generate, this thesis uses a fixed subset of 1,000, split 80/10/10 at the scenario level. That gives roughly 3.16 million test predictions.")


# ──────────────────── SLIDE 5 — Architecture ─────────────────
s = prs.slides.add_slide(BLANK)
add_title_bar(s, 5, "Base surrogate: PointNetTransfGAT (Elena's architecture)")

add_image_safe(s, Inches(0.5), Inches(1.4), Inches(8.5), Inches(4.7),
               "slide_architecture.png")

add_bullets(s, Inches(9.2), Inches(1.5), Inches(3.9), Inches(5), [
    "Three GNN ideas:",
    "  PointNetConv ×2",
    "  TransformerConv ×2",
    "  GATConv ×2",
    "",
    "≈ 1.42 M trainable",
    "parameters",
    "",
    "Dropout active in:",
    "  PointNet MLPs",
    "  between Transformers",
    "",
    "(Adapted from",
    "Natterer et al. 2025)",
], size=13, color=TEXT_DARK)

add_footer(s, 5)
set_notes(s, "The surrogate I build on is PointNetTransfGAT from Elena's preprint. It combines three GNN ideas. PointNetConv handles local geometry. TransformerConv reaches longer dependencies. GATConv handles final node-level mixing. About 1.4 million trainable parameters. Dropout in PointNet MLPs and between TransformerConv layers — that's what makes MC Dropout possible later. I did not change this architecture; I take Elena's design as given. What I vary is the training configuration.")


# ──────────────────── SLIDE 6 — Base trials T1–T8 ────────────
s = prs.slides.add_slide(BLANK)
add_title_bar(s, 6, "Base trials T1–T8 and the T8 UQ checkpoint")

# table-like structure (use a real PowerPoint table)
rows_t6 = [
    ["Trial", "Final layer", "Dropout", "R²", "MC-Dropout?"],
    ["T1", "Linear", "OFF", "0.7860", "NO — every pass identical"],
    ["T3, T4", "GATConv", "ON", "0.22 – 0.24", "weighted-MSE failures"],
    ["T8", "GATConv", "ON (p=0.2)", "0.5957", "YES — primary UQ base"],
]
table_shape = s.shapes.add_table(len(rows_t6), 5,
                                  Inches(0.7), Inches(1.5),
                                  Inches(12.0), Inches(2.5)).table
for r_idx, row in enumerate(rows_t6):
    for c_idx, val in enumerate(row):
        cell = table_shape.cell(r_idx, c_idx)
        cell.text = val
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(13)
                run.font.name = "Calibri"
                run.font.color.rgb = TEXT_DARK
                if r_idx == 0:
                    run.font.bold = True
                    run.font.color.rgb = TUM_BLUE_DARK
        cell.fill.solid()
        if r_idx == 0:
            cell.fill.fore_color.rgb = TUM_BLUE_PALE
        elif r_idx == 3:        # T8 row highlighted
            cell.fill.fore_color.rgb = TUM_GREY_LIGHT
        else:
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

add_bullets(s, Inches(0.7), Inches(4.4), Inches(12.0), Inches(2.3), [
    "→ T1 has highest absolute R², but dropout is OFF → MC Dropout undefined",
    "→ T3, T4 use weighted-MSE → fail to generalise",
    "→ T8 = strongest UQ-compatible base (NOT highest R² in project)",
    "",
    "Honest caveats",
    "  Selection used test R² (acknowledged as exploratory)",
    "  T8 sits below every Deep Ensemble member [0.640, 0.650]",
], size=15)

add_footer(s, 6)
set_notes(s, "I trained eight base configurations of the same architecture, varying dropout, batch size, learning rate, and loss. T1 has the highest R² at 0.79 but its dropout flag is off, so MC Dropout reduces to a constant — every pass produces an identical output. T1 is not usable for UQ work. T3 and T4 used weighted-MSE and didn't generalise. Among the dropout-enabled GATConv trials, T8 achieves the highest test R² at 0.5957. So T8 is the primary UQ checkpoint. Two honest caveats: T8 is the strongest UQ-compatible base, not the strongest predictor overall, and selection used test performance — acknowledged as exploratory.")


# ──────────────────── SLIDE 7 — Why UQ ───────────────────────
s = prs.slides.add_slide(BLANK)
add_title_bar(s, 7, "Why uncertainty quantification is needed")

add_text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.7),
         "A point prediction alone does not tell a planner WHEN to trust it.",
         size=20, bold=True, color=TUM_BLUE_DARK)

add_bullets(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.5), [
    "Training set covers only 1,000 of ≈ 10,000 possible scenarios",
    "Out-of-distribution scenarios degrade silently — confident, but wrong",
    "Infrastructure decisions need PER-PREDICTION reliability — not just average accuracy",
], size=17)

# goal box
goal = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(5.3),
                          Inches(11.5), Inches(1.3))
goal.fill.solid(); goal.fill.fore_color.rgb = TUM_BLUE_PALE
goal.line.color.rgb = ACCENT_LINE
tf = goal.text_frame
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2)
tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "GOAL  —  per-segment uncertainty signal that flags risky predictions for MATSim follow-up"
r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = TUM_BLUE_DARK
r.font.name = "Calibri"

add_footer(s, 7)
set_notes(s, "So far the surrogate gives one number per segment — a predicted delta v. That doesn't tell a planner whether to trust the number for any specific segment. Two issues. First, the training set covers only 1,000 of 10,000 scenarios, so a planner will eventually feed in something the model has never seen. Second, when that happens, the model doesn't say 'I'm not sure' — it just outputs a number. The goal of the UQ work is to add a per-segment reliability signal — a number that says 'this one's probably fine' or 'this one needs MATSim verification'.")


# ──────────────────── SLIDE 8 — MC Dropout method ────────────
s = prs.slides.add_slide(BLANK)
add_title_bar(s, 8, "MC Dropout method and the S = 30 choice")

add_bullets(s, Inches(0.7), Inches(1.4), Inches(6.0), Inches(5), [
    "Keep dropout ACTIVE at inference",
    "→ S stochastic forward passes per input",
    "",
    "For each segment:",
    "  prediction = mean over S passes",
    "  σ (uncertainty) = std-dev over S passes",
    "",
    "S = 30 chosen via convergence sweep",
    "on 10 test graphs",
    "  ρ at S = 30 / S = 50 = 0.466 / 0.471 (+1.03 %)",
    "  Final UQ eval uses ALL 100 test graphs at S = 30",
    "  ≈ 228 min inference",
], size=14)

add_image_safe(s, Inches(7.0), Inches(1.4), Inches(6.0), Inches(5.0),
               "fig05_s_convergence.png")

add_footer(s, 8)
set_notes(s, "The first UQ method is Monte Carlo Dropout. We keep dropout active at inference. For each input segment we run S stochastic forward passes — each pass drops a different random subset of activations — and we get S predictions. The mean is the point prediction; the standard deviation is sigma. S = 30 was chosen via a convergence sweep on ten test graphs because the full 100-graph sweep across many S values would have been too expensive on a T4. The Spearman correlation plateaus around S = 25, and going to 50 adds only one percent. So S = 30 sits on the plateau, and final UQ numbers are computed once on the full 100-graph test set at S = 30 — about 228 minutes of inference.")


# ──────────────────── SLIDE 9 — MC Dropout results ───────────
s = prs.slides.add_slide(BLANK)
add_title_bar(s, 9, "MC Dropout results: ranking signal, not calibrated probability")

# small table
rows_t9 = [
    ["Metric", "Value", "Reading"],
    ["Spearman ρ pooled", "0.4820", "useful ranking"],
    ["Mean per-graph ρ (95 % CI)", "0.464 [0.460, 0.469]", "stable across graphs"],
    ["k₉₅ (vs Gaussian 1.96)", "11.66", "σ severely under-dispersed"],
    ["Coverage at 1.96σ", "54.85 %", "nominal 95 %"],
]
t = s.shapes.add_table(len(rows_t9), 3,
                       Inches(0.7), Inches(1.5),
                       Inches(7.0), Inches(2.7)).table
for r_idx, row in enumerate(rows_t9):
    for c_idx, val in enumerate(row):
        cell = t.cell(r_idx, c_idx)
        cell.text = val
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(12)
                run.font.name = "Calibri"
                run.font.color.rgb = TEXT_DARK
                if r_idx == 0:
                    run.font.bold = True
                    run.font.color.rgb = TUM_BLUE_DARK
        cell.fill.solid()
        cell.fill.fore_color.rgb = TUM_BLUE_PALE if r_idx == 0 else RGBColor(0xFF, 0xFF, 0xFF)

add_bullets(s, Inches(0.7), Inches(4.4), Inches(7.0), Inches(2.5), [
    "→ σ orders errors well — useful as a ranking signal",
    "→ σ is NOT a calibrated standard deviation",
    "→ Calibration cascade on the next slide closes this gap",
], size=15, color=TUM_BLUE_DARK)

add_image_safe(s, Inches(7.9), Inches(1.4), Inches(5.2), Inches(5.0),
               "fig11_k95_comparison.png")

add_footer(s, 9)
set_notes(s, "These are the headline MC Dropout numbers on T8. Pooled Spearman correlation between sigma and absolute error is 0.482; per-graph mean is 0.464 with a tight bootstrap CI. Sigma orders errors reasonably well; no graph drops below 0.41. But the absolute scale is a different story. K95 = 11.66 against the ideal Gaussian 1.96. If we use the standard 1.96σ interval, only 55% of points are covered, against a nominal 95%. So MC Dropout works as a ranking signal but raw sigma is not a calibrated probability.")


# ──────────────────── SLIDE 10 — Calibration ─────────────────
s = prs.slides.add_slide(BLANK)
add_title_bar(s, 10, "Calibration: temperature scaling + conformal prediction")

add_bullets(s, Inches(0.6), Inches(1.4), Inches(6.5), Inches(5.5), [
    "Two complementary post-hoc fixes — no retraining needed",
    "",
    "TEMPERATURE SCALING — single scalar T",
    "  T★ = 2.887",
    "  ECE: 0.356 → 0.034   (−90.5 %)",
    "  1σ coverage: 32.7 % → 68.0 %",
    "  k₉₅: 11.66 → 4.04",
    "",
    "SPLIT + ADAPTIVE CONFORMAL",
    "  q₉₀ = 9.92,  PICP₉₀ = 90.02 %",
    "  q₉₅ = 14.68, PICP₉₅ = 95.01 %",
    "  Adaptive narrows decile range:",
    "    [59 %, 98 %] → [83.7 %, 96.4 %]",
], size=14)

add_image_safe(s, Inches(7.4), Inches(1.4), Inches(5.7), Inches(5.0),
               "fig15_conditional_coverage_by_decile.png")

add_footer(s, 10)
set_notes(s, "Two post-hoc fixes that don't require retraining. First, temperature scaling — a single scalar T learned by minimising ECE. Optimal T = 2.89. ECE drops from 0.356 to 0.034 — about 90 percent. 1σ coverage goes from 33% to 68% — almost exactly the Gaussian target. Second, conformal prediction. Standard split conformal hits PICP = 90.02% and 95.01% — bullseye on the nominal targets. Adaptive conformal divides the residual by sigma so the interval scales with uncertainty. Marginal coverage stays the same; conditional coverage across deciles tightens dramatically — from a 59 to 98 percent spread under standard, down to 83 to 96 under adaptive. Caveat: the conformal guarantee rests on exchangeability, enforced at scenario level but only approximated at node level.")


# ──────────────────── SLIDE 11 — Stratified UQ ───────────────
s = prs.slides.add_slide(BLANK)
add_title_bar(s, 11, "Stratified UQ: where the method is weakest")

add_image_safe(s, Inches(0.5), Inches(1.4), Inches(7.5), Inches(5.0),
               "fig28_stratified_uq_quartiles.png")

# quartile mini-table on right
rows_t11 = [
    ["Q", "mean |Δv|", "ρ", "MAE"],
    ["Q1", "0.00", "0.721 *", "1.24"],
    ["Q2", "0.45", "0.529",  "1.67"],
    ["Q3", "2.46", "0.162",  "2.80"],
    ["Q4", "14.14", "0.100", "10.08"],
]
t = s.shapes.add_table(len(rows_t11), 4,
                       Inches(8.3), Inches(1.5),
                       Inches(4.7), Inches(2.6)).table
for r_idx, row in enumerate(rows_t11):
    for c_idx, val in enumerate(row):
        cell = t.cell(r_idx, c_idx)
        cell.text = val
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(13)
                run.font.name = "Calibri"
                run.font.color.rgb = TEXT_DARK
                if r_idx == 0:
                    run.font.bold = True
                    run.font.color.rgb = TUM_BLUE_DARK
        cell.fill.solid()
        if r_idx == 0:
            cell.fill.fore_color.rgb = TUM_BLUE_PALE
        elif r_idx == 4:
            cell.fill.fore_color.rgb = TUM_CORAL
        else:
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

add_bullets(s, Inches(8.3), Inches(4.3), Inches(4.7), Inches(2.5), [
    "* Q1 ρ partly mechanical",
    "  (zero-effect segments)",
    "",
    "→ Method weakest where",
    "  policy effect is largest",
    "",
    "→ Mitigation: route high σ +",
    "  high |Δv̂| back to MATSim",
], size=12)

add_footer(s, 11)
set_notes(s, "This is the most policy-relevant slide. We slice the test set by true delta v into four quartiles. Q1 has mean delta v zero — segments with no policy effect. Q4 has mean fourteen, max around 230 — segments where policy moves traffic most. Within each quartile we recompute Spearman rho. Q1 is 0.72 but partly mechanical because Q1 is dominated by zero-effect segments. Real story: collapse between Q2 and Q3 from 0.53 to 0.16. By Q4 it's 0.10 — basically uninformative. MAE shows the same — gradual on Q1 to Q3, then a jump to ten on Q4. Sigma rises with difficulty but only doubles, while difficulty rises 14×. The deployment framework explicitly handles this — high sigma plus high predicted delta v gets routed to MATSim.")


# ──────────────────── SLIDE 12 — DE vs MC Dropout ────────────
s = prs.slides.add_slide(BLANK)
add_title_bar(s, 12, "Ensembles: MC Dropout vs Deep Ensemble")

rows_t12 = [
    ["Metric",          "T8 (MC Dropout)", "Deep Ensemble"],
    ["R² (point)",      "0.5957",          "0.6841 (+14.8 %)"],
    ["MAE (veh/h)",     "3.957",           "3.485"],
    ["Spearman ρ (UQ)", "0.4820",          "0.3997 (−17 %)"],
    ["k₉₅ (calibration)","11.66",           "15.18"],
]
t = s.shapes.add_table(len(rows_t12), 3,
                       Inches(0.7), Inches(1.5),
                       Inches(8.0), Inches(2.7)).table
for r_idx, row in enumerate(rows_t12):
    for c_idx, val in enumerate(row):
        cell = t.cell(r_idx, c_idx)
        cell.text = val
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(13)
                run.font.name = "Calibri"
                run.font.color.rgb = TEXT_DARK
                if r_idx == 0:
                    run.font.bold = True
                    run.font.color.rgb = TUM_BLUE_DARK
        cell.fill.solid()
        cell.fill.fore_color.rgb = TUM_BLUE_PALE if r_idx == 0 else RGBColor(0xFF, 0xFF, 0xFF)

add_bullets(s, Inches(0.7), Inches(4.4), Inches(8.0), Inches(2.6), [
    "5 members, T8 hyperparams, seeds {42, 137, 256, 389, 512}",
    "",
    "→ Deep Ensemble wins POINT accuracy",
    "→ MC Dropout wins UQ ranking + calibration",
    "→ Complementary, not competing",
    "→ Pair both at ~5× compute for high-stakes use",
    "",
    "Note: DE = strongest GNN-based predictor.",
    "      XGBoost (0.7414) is stronger overall.",
], size=14)

add_image_safe(s, Inches(8.9), Inches(1.5), Inches(4.2), Inches(4.5),
               "fig26_deep_ensemble_member_r2.png")

add_footer(s, 12)
set_notes(s, "How does MC Dropout compare with the standard alternative — a Deep Ensemble of independently trained networks? I trained five members with T8 hyperparameters, different seeds. Ensemble prediction is the mean; disagreement gives uncertainty. Deep Ensemble is the strongest GNN-based point predictor — R² 0.68, almost 15% above T8. Individual members all in [0.640, 0.650]. Averaging lifts the mean above every member. But on uncertainty quality, DE is actually weaker — Spearman 0.40 vs MC Dropout's 0.48; k95 = 15 vs 11.66. So DE wins point accuracy, MC Dropout wins UQ ranking. Complementary, not competing. Run both at ~5× training compute for high-stakes use.")


# ──────────────────── SLIDE 13 — T9 / T10 / T11 ──────────────
s = prs.slides.add_slide(BLANK)
add_title_bar(s, 13, "Uncertainty-aware training: T9, T10, T11")

rows_t13 = [
    ["Trial", "Head + loss",          "Backbone",                 "R²",      "Gates"],
    ["T9",    "heteroscedastic (NLL)", "frozen",                  "0.4991",  "2 / 3"],
    ["T10",   "CQR (pinball loss)",   "UNFROZEN  (lr 5e-4)",      "0.4057",  "3 / 6  ✗"],
    ["T11",   "CQR (same, same lr)",  "FROZEN",                   "0.5835",  "6 / 6  ✓"],
]
t = s.shapes.add_table(len(rows_t13), 5,
                       Inches(0.6), Inches(1.4),
                       Inches(8.5), Inches(2.4)).table
for r_idx, row in enumerate(rows_t13):
    for c_idx, val in enumerate(row):
        cell = t.cell(r_idx, c_idx)
        cell.text = val
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(12)
                run.font.name = "Calibri"
                run.font.color.rgb = TEXT_DARK
                if r_idx == 0:
                    run.font.bold = True
                    run.font.color.rgb = TUM_BLUE_DARK
        cell.fill.solid()
        if r_idx == 0:
            cell.fill.fore_color.rgb = TUM_BLUE_PALE
        elif r_idx == 2:           # T10 row — fail
            cell.fill.fore_color.rgb = TUM_CORAL
        elif r_idx == 3:           # T11 row — pass
            cell.fill.fore_color.rgb = TUM_BLUE_LIGHT
        else:
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

add_bullets(s, Inches(0.6), Inches(4.0), Inches(8.5), Inches(3), [
    "→ T10 vs T11 = clean SINGLE-KNOB ablation",
    "    same head, same loss, same data, SAME learning rate",
    "    only difference: backbone-trainability flag",
    "",
    "→ Empirical observation: freezing the MSE-trained backbone is",
    "    the strongest observed design choice in this setting",
    "→ NOT claimed as universal principle",
], size=14)

add_image_safe(s, Inches(9.3), Inches(1.4), Inches(3.8), Inches(5.0),
               "fig22_cqr_r2_progression.png")

add_footer(s, 13)
set_notes(s, "Three uncertainty-aware variants on top of T8's MSE-trained backbone. T9 freezes the backbone, trains a 134-parameter heteroscedastic head under NLL. R² drops to 0.50 missing the gate, but k95 reaches 2.84 — almost ideal. Partial positive. T10 takes the same backbone and bolts on a CQR head with pinball loss, fine-tunes the full backbone end-to-end at lr 5e-4. R² collapses to 0.41. T11 is identical to T10 — same head, same loss, same data, same learning rate — except the backbone is frozen. R² recovers to 0.58. T10 vs T11 is a clean single-knob ablation. The only design difference is the freezing flag. Pinball gradients reaching the full backbone reshape representations originally learned for MSE — that's the Seitzer mechanism. Freezing prevents that reshaping. Reported as empirical observation specific to this setting, not as universal principle.")


# ──────────────────── SLIDE 14 — Non-GNN baselines ───────────
s = prs.slides.add_slide(BLANK)
add_title_bar(s, 14, "Non-GNN baselines, limitations, threats to validity")

rows_t14 = [
    ["Model",              "R²",      "MAE",   "RMSE",   "UQ stack?"],
    ["XGBoost",            "0.7414",  "2.774", "5.693",  "none"],
    ["Deep Ensemble (GNN)","0.6841",  "3.485", "6.293",  "via members"],
    ["Random Forest",      "0.6612",  "3.263", "6.516",  "split conformal"],
    ["T8 (GNN, primary UQ)","0.5957", "3.957", "7.118",  "FULL"],
    ["MLP (sklearn-style)","0.4928",  "3.883", "7.973",  "none"],
]
t = s.shapes.add_table(len(rows_t14), 5,
                       Inches(0.5), Inches(1.4),
                       Inches(8.5), Inches(3.0)).table
for r_idx, row in enumerate(rows_t14):
    for c_idx, val in enumerate(row):
        cell = t.cell(r_idx, c_idx)
        cell.text = val
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(11)
                run.font.name = "Calibri"
                run.font.color.rgb = TEXT_DARK
                if r_idx == 0:
                    run.font.bold = True
                    run.font.color.rgb = TUM_BLUE_DARK
                elif r_idx == 1:   # XGBoost winner
                    run.font.bold = True
        cell.fill.solid()
        if r_idx == 0:
            cell.fill.fore_color.rgb = TUM_BLUE_PALE
        elif r_idx == 1:
            cell.fill.fore_color.rgb = TUM_ORANGE
        elif r_idx == 4:           # T8 highlighted
            cell.fill.fore_color.rgb = TUM_GREY_LIGHT
        else:
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

add_bullets(s, Inches(0.5), Inches(4.6), Inches(8.5), Inches(2.5), [
    "→ XGBoost wins POINT accuracy on this sparse target (88.7 % zero-mass)",
    "→ GNN's contribution is the per-prediction UQ pipeline trees lack in default form",
], size=14, color=TUM_BLUE_DARK)

# threats to validity panel
panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                           Inches(9.2), Inches(1.4),
                           Inches(3.9), Inches(5.4))
panel.fill.solid(); panel.fill.fore_color.rgb = TUM_BLUE_PALE
panel.line.color.rgb = ACCENT_LINE
add_text(s, Inches(9.4), Inches(1.55), Inches(3.6), Inches(0.5),
         "Threats to validity",
         size=15, bold=True, color=TUM_BLUE_DARK)
add_bullets(s, Inches(9.4), Inches(2.05), Inches(3.6), Inches(4.7), [
    "1. Dataset",
    "   1,000 of 10,000 scenarios",
    "",
    "2. External validity",
    "   Paris + capacity reduction only",
    "",
    "3. Model family",
    "   mainly PointNetTransfGAT",
    "",
    "4. Baseline scope",
    "   trees beat T8 on R²",
    "",
    "5. Metric limitation",
    "   strong UQ ≠ guaranteed",
    "   policy reliability",
], size=10, color=TEXT_DARK)

add_footer(s, 14)
set_notes(s, "Two honest moves. First, the non-GNN baselines. Dominik and Elena specifically asked for these. XGBoost is the strongest point predictor in the project at R² 0.74, beating T8 and the Deep Ensemble. Random Forest also beats T8. Consistent with literature on tree ensembles for tabular sparse-response targets. So the GNN's contribution is orthogonal — the per-prediction UQ pipeline trees don't provide in their default form. Second, threats to validity. Five honestly: only 1,000 of 10,000 scenarios; only Paris and capacity reductions; mainly one architecture family; the baseline gap I just discussed; and finally that strong UQ metrics don't on their own imply policy-level reliability.")


# ──────────────────── SLIDE 15 — Final takeaways ─────────────
s = prs.slides.add_slide(BLANK)
add_title_bar(s, 15, "Final takeaways")

add_bullets(s, Inches(0.7), Inches(1.4), Inches(12.0), Inches(4.5), [
    "1.  Calibrated UQ for a GNN traffic surrogate is achievable at tractable cost",
    "    — built as a 5-layer stack on T8.",
    "",
    "2.  MC Dropout works as a RANKING signal (ρ = 0.482),",
    "    but raw σ is NOT a calibrated probability (k₉₅ = 11.66).",
    "",
    "3.  Calibration cascade closes the gap:",
    "    temperature scaling → ECE 0.356 → 0.034",
    "    adaptive conformal → conditional coverage [59 %, 98 %] → [83.7 %, 96.4 %]",
    "",
    "4.  Freezing the MSE-trained backbone is the strongest observed design choice",
    "    when adding an uncertainty-aware head (T10 vs T11 single-knob ablation).",
    "",
    "5.  Method weakest where it matters most: Q4 → ρ = 0.10 → routed to MATSim.",
    "",
    "6.  XGBoost wins POINT accuracy. GNN wins per-prediction UQ.",
], size=14)

# scope + thanks panel
panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                           Inches(0.7), Inches(6.05), Inches(12.0), Inches(0.95))
panel.fill.solid(); panel.fill.fore_color.rgb = TUM_BLUE_PALE
panel.line.fill.background()
tf = panel.text_frame
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.15)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = ("SCOPE — one Paris network · one intervention · one GNN family · 1,000 of 10,000 scenarios   "
          "·   findings = empirical observations within this scope")
r.font.size = Pt(13); r.font.italic = True; r.font.color.rgb = TUM_BLUE_DARK
r.font.name = "Calibri"

# Thank you
add_text(s, Inches(0.7), Inches(7.05), Inches(12), Inches(0.4),
         "Thank you  —  happy to discuss",
         size=18, bold=True, color=TUM_BLUE_DARK, align=PP_ALIGN.CENTER)

set_notes(s, "To close — six takeaways. First, calibrated UQ is achievable at tractable cost as a layered stack on T8. Second, raw MC Dropout sigma is a ranking signal, not a calibrated probability. Third, the calibration cascade — temperature scaling and adaptive conformal — closes the gap. Fourth, freezing the MSE-trained backbone is the strongest observed design choice for uncertainty-aware heads; T10 vs T11 single-knob ablation. Fifth, the method is weakest exactly where the policy effect is largest — Q4 — handled by routing to MATSim. Sixth, XGBoost wins point accuracy; the GNN wins per-prediction UQ. Everything is empirical and scoped. Future work: full corpus replication, tree-native UQ, architecturally diverse ensembles, multi-city validation. Thank you — happy to discuss.")


# ───────────────────────── Save ──────────────────────────────
prs.save(OUT_PATH)
print(f"Saved deck → {OUT_PATH}")
print(f"15 slides, {os.path.getsize(OUT_PATH) / 1024:.1f} KB")
